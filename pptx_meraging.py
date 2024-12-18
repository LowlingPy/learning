from pptx import Presentation

# List of PowerPoint files to merge
files = ["1-5.pptx", "6-13.pptx", "14-15.pptx", "16-21.pptx", "22-24", "25-27",
         "28-30.pptx", "31-37.pptx", "38-40.pptx", "41-44.pptx", "45-46.pptx",
         '47-50.pptx']

# Create a new presentation
merged_ppt = Presentation()

for file in files:
    ppt = Presentation(file)
    for slide in ppt.slides:
        slide_copy = merged_ppt.slides.add_slide(merged_ppt.slide_layouts[0])
        for shape in slide.shapes:
            if shape.text:
                slide_copy.shapes.title.text = shape.text

# Save the merged presentation
merged_ppt.save("1-50.pptx")
